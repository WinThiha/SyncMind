from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x02, 0x06, 0x17)
NAVY       = RGBColor(0x0f, 0x17, 0x2a)
CARD       = RGBColor(0x0d, 0x1f, 0x3c)
CARD_AI    = RGBColor(0x0d, 0x2a, 0x4a)
BLUE       = RGBColor(0x3b, 0x82, 0xf6)
BLUE_MID   = RGBColor(0x25, 0x63, 0xeb)
BLUE_LIGHT = RGBColor(0x60, 0xa5, 0xfa)
BLUE_PALE  = RGBColor(0x93, 0xc5, 0xfd)
WHITE      = RGBColor(0xf8, 0xfa, 0xfc)
MUTED      = RGBColor(0x94, 0xa3, 0xb8)
BORDER     = RGBColor(0x1e, 0x3a, 0x5f)

# ── Helpers ───────────────────────────────────────────────────────────────────
def inches(n): return Inches(n)
def pt(n):     return Pt(n)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, radius=False):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1 if not radius else 5,
        inches(x), inches(y), inches(w), inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name='Calibri', wrap=True):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)

def card(slide, x, y, w, h, fill=None, border=None):
    f = fill or CARD
    b = border or BORDER
    add_rect(slide, x, y, w, h, f, b)

def accent_bar(slide, y=0.38, h=0.44):
    add_rect(slide, 0.45, y, 0.05, h, BLUE)

def slide_title(slide, text, sub=None):
    accent_bar(slide)
    add_text(slide, text, 0.62, 0.30, 12.0, 0.55,
             font_size=22, bold=True, color=WHITE)
    if sub:
        add_text(slide, sub, 0.62, 0.82, 12.0, 0.32,
                 font_size=10, color=MUTED, italic=True)
    add_rect(slide, 0.45, 1.08, 12.4, 0.02, BORDER)

def divider_v(slide, x, y, h):
    add_rect(slide, x, y, 0.02, h, BORDER)

def page_num(slide, n, total=20):
    add_text(slide, f'{n} / {total}', 11.9, 7.1, 1.2, 0.28,
             font_size=8, color=BORDER, align=PP_ALIGN.RIGHT)

def bullet_list(slide, items, x, y, w, lh=0.48, fs=12):
    for i, item in enumerate(items):
        if item == '':
            continue
        is_section = item.startswith('##')
        text = item.replace('## ', '')
        color = BLUE_LIGHT if is_section else WHITE
        size  = fs - 1    if is_section else fs
        prefix = ''        if is_section else '  -  '
        add_text(slide, prefix + text, x, y + i * lh, w, lh,
                 font_size=size, bold=is_section, color=color)

def top_bar(slide, x, y, w, color=BLUE):
    add_rect(slide, x, y, w, 0.055, color)

# =============================================================================
# SLIDE 1  Title
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
add_rect(s, 9.8, 0, 3.53, 7.5, NAVY)
add_rect(s, 9.8, 0, 0.06, 7.5, BLUE)
add_text(s, 'SyncMind', 0.6, 1.5, 8.8, 1.5,
         font_size=72, bold=True, color=WHITE)
add_text(s, 'AI-Powered Project and Issue Management', 0.6, 3.1, 8.8, 0.6,
         font_size=20, color=BLUE_LIGHT)
add_rect(s, 0.6, 3.9, 3.8, 0.05, BLUE)
add_text(s, 'Smarter teamwork. Less manual effort.', 0.6, 4.1, 8.8, 0.45,
         font_size=13, italic=True, color=MUTED)
add_text(s, 'Version 1.0  |  2026', 10.1, 3.55, 3.0, 0.45,
         font_size=13, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
page_num(s, 1)

# =============================================================================
# SLIDE 2  The Problem
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'The Problem', 'Everyday pain points that slow development teams down')

problems = [
    ('01', 'Writing the same issue descriptions over and over from scratch'),
    ('02', 'Creating duplicate tasks because search did not surface the existing one'),
    ('03', 'Guessing who should own a task without clear reasoning'),
    ('04', 'Reading through dozens of comments just to find one key decision'),
    ('05', 'Losing sight of whether a deadline is still on track'),
]
for i, (num, text) in enumerate(problems):
    card(s, 0.45, 1.22 + i * 1.02, 12.2, 0.88)
    add_rect(s, 0.45, 1.22 + i * 1.02, 0.05, 0.88, BLUE)
    add_text(s, num, 0.62, 1.36 + i * 1.02, 0.55, 0.44,
             font_size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    divider_v(s, 1.22, 1.32 + i * 1.02, 0.55)
    add_text(s, text, 1.4, 1.36 + i * 1.02, 11.0, 0.52, font_size=13)
page_num(s, 2)

# =============================================================================
# SLIDE 3  What is SyncMind
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'What is SyncMind?')

card(s, 0.45, 1.22, 12.2, 1.55)
add_text(s, 'SyncMind is a web-based project management tool with AI built into the workflow, not added on top.',
         0.7, 1.34, 11.8, 0.5, font_size=15, bold=True)
add_text(s, 'It gives teams a structured space to create projects, manage issues, track milestones, and collaborate — while AI handles the writing, searching, assigning, and summarising automatically.',
         0.7, 1.86, 11.8, 0.62, font_size=12, color=MUTED)

# 5 pillars — cw=2.3, gap=0.19 → total span = 5*2.3 + 4*0.19 = 12.26
pillars = [
    ('Web App',       'Next.js + Laravel'),
    ('Team-Based',    'Roles and Permissions'),
    ('AI-Native',     '5 AI Capabilities'),
    ('API-First',     'Mobile-Ready Backend'),
    ('Multi-Lingual', '5 Languages'),
]
cw, gap = 2.3, 0.19
for i, (label, sub) in enumerate(pillars):
    cx = 0.45 + i * (cw + gap)
    card(s, cx, 3.05, cw, 2.0)
    top_bar(s, cx, 3.05, cw)
    add_text(s, label, cx + 0.1, 3.28, cw - 0.2, 0.45,
             font_size=13, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, sub,   cx + 0.1, 3.82, cw - 0.2, 0.38,
             font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, 'The goal is not to change how teams work — it is to remove the friction so they can focus on the actual work.',
         0.45, 5.28, 12.2, 0.42, font_size=12, italic=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
page_num(s, 3)

# =============================================================================
# SLIDE 4  Feature Map
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Feature Map', 'Everything SyncMind covers')

# 8 rows — step=0.73, card height=0.63 → last card ends at 1.18+7*0.73+0.63=6.92
rows = [
    ('Authentication', 'Email registration, Google OAuth, password reset, email verification, session management', False),
    ('Projects',       'Create and configure projects, member management, roles, email invitations', False),
    ('Issues',         'Create, edit, filter, search, track, comment, and audit every task', False),
    ('Milestones',     'Delivery goals with deadlines, automatic progress tracking and overdue detection', False),
    ('Collaboration',  'Markdown comments, field-level change history, unified activity timeline', False),
    ('AI',             'Auto-fill  |  Duplicate detection  |  Assignee AI  |  Semantic search  |  Thread summarisation', True),
    ('Localization',   'UI available in 5 languages — English, Myanmar, Japanese, Vietnamese, Khmer', False),
    ('App Shell',      'Dark mode, collapsible sidebar, glassmorphism design, responsive layout', False),
]
for i, (area, items, is_ai) in enumerate(rows):
    fill  = CARD_AI if is_ai else CARD
    bord  = BLUE    if is_ai else BORDER
    card(s, 0.45, 1.18 + i * 0.73, 12.2, 0.63, fill=fill, border=bord)
    if is_ai:
        add_rect(s, 0.45, 1.18 + i * 0.73, 0.05, 0.63, BLUE)
    add_text(s, area, 0.65, 1.26 + i * 0.73, 2.55, 0.38,
             font_size=11, bold=is_ai, color=BLUE_LIGHT)
    divider_v(s, 3.28, 1.28 + i * 0.73, 0.38)
    add_text(s, items, 3.45, 1.26 + i * 0.73, 9.0, 0.38,
             font_size=10, color=WHITE if is_ai else MUTED)
page_num(s, 4)

# =============================================================================
# SLIDE 5  Localization
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Localization — Multi-Language Support',
            'SyncMind is available in five languages out of the box')

card(s, 0.45, 1.18, 12.2, 1.05)
add_text(s, 'Every UI label, error message, button, and tooltip is translated. Users can switch their interface language from their profile settings at any time, with no page reload required.',
         0.65, 1.3, 11.8, 0.82, font_size=12, color=MUTED)

# 5 language cards — cw=2.28, gap=0.235 → 5*2.28+4*0.235=12.34
langs = [
    ('EN',  'English',    'Base language'),
    ('MM',  'Myanmar',    u'မြန်မာဗာသာ'),
    ('JP',  'Japanese',   u'日本語'),
    ('VI',  'Vietnamese', u'Tiếng Việt'),
    ('KM',  'Khmer',      u'ភាសាខ្មែរ'),
]
cw, gap = 2.28, 0.235
cy, ch  = 2.45, 3.2
for i, (code, name, native) in enumerate(langs):
    cx = 0.45 + i * (cw + gap)
    card(s, cx, cy, cw, ch, fill=CARD_AI, border=BLUE)
    top_bar(s, cx, cy, cw)
    add_text(s, code,   cx + 0.1, cy + 0.16, cw - 0.2, 0.7,
             font_size=32, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(s, cx + 0.15, cy + 0.98, cw - 0.3, 0.025, BORDER)
    add_text(s, name,   cx + 0.1, cy + 1.12, cw - 0.2, 0.42,
             font_size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, native, cx + 0.1, cy + 1.65, cw - 0.2, 0.42,
             font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

card(s, 0.45, 5.85, 12.2, 0.72)
add_text(s, 'Implementation: structured JSON translation files — one file per language. Adding a new language requires only a new translation file with no code changes.',
         0.65, 5.96, 11.8, 0.48, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
page_num(s, 5)

# =============================================================================
# SLIDE 6  Project and Team Management
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Project and Team Management')

card(s, 0.45, 1.18, 5.95, 5.85)
add_text(s, 'Project Setup', 0.65, 1.3, 5.6, 0.4,
         font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'Create with name, key (issue prefix), icon, and issue types',
    'Project key used as issue prefix  e.g.  SYNC-1, SYNC-2',
    'Edit name and issue types from project settings anytime',
], 0.65, 1.82, 5.6, lh=0.6, fs=12)
add_rect(s, 0.65, 3.68, 5.4, 0.02, BORDER)
add_text(s, 'Ownership and Danger Zone', 0.65, 3.82, 5.6, 0.4,
         font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'Transfer ownership to another admin',
    'Delete project  (creator only)',
], 0.65, 4.32, 5.6, lh=0.6, fs=12)

card(s, 6.75, 1.18, 5.95, 5.85)
add_text(s, 'Membership', 6.95, 1.3, 5.6, 0.4,
         font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'Invite members by email — one-click acceptance link',
    'Admin: manage members and issues',
    'Normal: view and participate in collaboration',
    'Change roles or remove members at any time',
    'Members see only projects they belong to',
], 6.95, 1.82, 5.6, lh=0.6, fs=12)
page_num(s, 6)

# =============================================================================
# SLIDE 7  Issues and Milestones
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Issues and Milestones')

card(s, 0.45, 1.18, 5.95, 5.85)
top_bar(s, 0.45, 1.18, 5.95)
add_text(s, 'Issues', 0.65, 1.36, 5.5, 0.44, font_size=14, bold=True, color=BLUE_LIGHT)
add_text(s, 'Day-to-day tasks your team works on.', 0.65, 1.78, 5.5, 0.32,
         font_size=10, italic=True, color=MUTED)
bullet_list(s, [
    'Summary, description, type, priority, status',
    'Assignee, estimated and actual hours',
    'Auto-generated keys  KEY-1, KEY-2 ...',
    'Full field-level change history',
    'Comments with markdown support',
    'Slide-over panel for quick updates',
    'Soft delete  nothing lost permanently',
], 0.65, 2.18, 5.5, lh=0.52, fs=11)

card(s, 6.75, 1.18, 5.95, 5.85)
top_bar(s, 6.75, 1.18, 5.95)
add_text(s, 'Milestones', 6.95, 1.36, 5.5, 0.44, font_size=14, bold=True, color=BLUE_LIGHT)
add_text(s, 'Delivery goals with automatic progress tracking.', 6.95, 1.78, 5.5, 0.32,
         font_size=10, italic=True, color=MUTED)
bullet_list(s, [
    'Name, description, start date, due date',
    'Status: Open, In Progress, Closed',
    'Assign issues via dropdown on the issue form',
    'Progress auto-calculates from issue statuses',
    'Overdue flag triggers when deadline passes',
    'Timeline view across all milestones',
    'Milestone detail lists all linked issues',
], 6.95, 2.18, 5.5, lh=0.52, fs=11)
page_num(s, 7)

# =============================================================================
# SLIDE 8  Demo
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, BLUE)
add_text(s, 'Live Demonstration', 0.5, 0.65, 12.3, 0.95,
         font_size=40, bold=True, color=WHITE)
add_rect(s, 0.5, 1.65, 5.0, 0.05, BLUE)

steps = [
    '01   Login and dashboard',
    '02   Create project, invite a member',
    '03   Create a milestone with a due date',
    '04   Create issue  -  trigger AI auto-fill',
    '05   Watch duplicate detection fire in real time',
    '06   Accept an AI assignee recommendation',
    '07   Link issue to milestone, see progress update',
    '08   Open an existing thread, trigger summarisation',
    '09   Toggle AI semantic search, query in natural language',
]
for i, step in enumerate(steps):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    add_text(s, step, 0.55 + col * 6.55, 1.92 + row * 0.9, 6.2, 0.72,
             font_size=14, color=WHITE)
    add_rect(s, 0.55 + col * 6.55, 2.62 + row * 0.9, 5.5, 0.015, BORDER)
page_num(s, 8)

# =============================================================================
# SLIDE 9  Architecture
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'System Architecture', 'Two independent applications communicating over HTTP')

boxes = [
    ('Browser',          'User Interface',                  0.5,  1.55, 2.4,  1.15),
    ('Next.js Frontend', 'TypeScript  |  React 19\nPort 3000', 3.35, 1.55, 2.85, 1.15),
    ('Laravel Backend',  'PHP 8.2  |  REST API\nPort 8000',    7.05, 1.55, 2.85, 1.15),
    ('PostgreSQL 16',    'Data and pgvector\nEmbeddings',      10.4, 1.55, 2.4,  1.15),
    ('OpenRouter',       'AI Model Gateway\nLLM Calls',        7.05, 3.3,  2.85, 1.05),
    ('Docker Compose',   'Orchestrates all services',         10.4, 3.3,  2.4,  1.05),
]
for label, sub, x, y, w, h in boxes:
    card(s, x, y, w, h)
    top_bar(s, x, y, w)
    add_text(s, label, x+0.1, y+0.15, w-0.2, 0.38,
             font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, sub, x+0.1, y+0.56, w-0.2, 0.52,
             font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Connectors
add_rect(s, 2.92, 2.1,  0.42, 0.025, BLUE)
add_rect(s, 6.22, 2.1,  0.82, 0.025, BLUE)
add_rect(s, 9.92, 2.1,  0.48, 0.025, BLUE)
add_rect(s, 8.46, 2.72, 0.025, 0.57, BLUE)

add_text(s, 'Axios  |  withCredentials  |  Sanctum Cookie Auth  |  CSRF Protected',
         0.45, 4.6, 12.2, 0.32, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)
card(s, 0.45, 5.05, 12.2, 0.88)
add_text(s, 'Everything runs in Docker. The backend is a pure REST API. This makes it straightforward to add new clients such as a mobile app without any changes to the server.',
         0.65, 5.2, 11.8, 0.58, font_size=11, align=PP_ALIGN.CENTER)
page_num(s, 9)

# =============================================================================
# SLIDE 10  Tech Stack
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Technology Stack')

rows = [
    ('Frontend',       'Next.js 16  |  TypeScript  |  React 19',      'Fast SPA, strong typing, React Compiler enabled'),
    ('Backend',        'Laravel 12  |  PHP 8.2',                       'Mature API framework, built-in auth and validation'),
    ('Database',       'PostgreSQL 16 + pgvector',                     'Relational data and vector similarity in one database'),
    ('Auth',           'Laravel Sanctum',                              'Secure cookie-based sessions with CSRF protection'),
    ('AI Gateway',     'OpenRouter',                                   'Flexible model routing with OpenAI-compatible API'),
    ('Infrastructure', 'Docker Compose',                               'Consistent environment, easy local and production setup'),
    ('Design System',  'Glassmorphism  |  Framer Motion  |  Lucide',   'Polished, animated, consistent UI across all surfaces'),
]

# Header
card(s, 0.45, 1.18, 12.2, 0.44, fill=NAVY)
for text, xpos, ww in [('LAYER', 0.65, 2.1), ('TECHNOLOGY', 2.95, 4.15), ('PURPOSE', 7.3, 5.1)]:
    add_text(s, text, xpos, 1.26, ww, 0.3,
             font_size=9, bold=True, color=MUTED)

for i, (layer, tech, why) in enumerate(rows):
    card(s, 0.45, 1.65 + i * 0.815, 12.2, 0.72)
    add_text(s, layer, 0.65, 1.74 + i * 0.815, 2.1,  0.42, font_size=11, bold=True, color=BLUE_LIGHT)
    divider_v(s, 2.82, 1.74 + i * 0.815, 0.42)
    add_text(s, tech,  2.98, 1.74 + i * 0.815, 4.15, 0.42, font_size=11)
    divider_v(s, 7.2,  1.74 + i * 0.815, 0.42)
    add_text(s, why,   7.35, 1.74 + i * 0.815, 5.1,  0.42, font_size=10, color=MUTED)
page_num(s, 10)

# =============================================================================
# SLIDE 11  AI Overview
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI in SyncMind', 'Five capabilities built into the daily workflow')

add_text(s, 'AI appears at five specific moments — each targeting a real pain point, not a showcase feature.',
         0.45, 1.12, 12.2, 0.38, font_size=12, color=MUTED)

items = [
    ('01', 'Auto-Fill',            'When writing an issue',     'Generates description, type, priority and hours from a short summary'),
    ('02', 'Duplicate Detection',  'When creating an issue',    'Real-time vector similarity check prevents creating work that already exists'),
    ('03', 'Assignee AI',          'When assigning an issue',   'Ranks team members by fit with a written reason for each recommendation'),
    ('04', 'Semantic Search',      'When searching issues',     'Finds issues by meaning — not just exact keyword matching'),
    ('05', 'Thread Summarisation', 'When reviewing discussion', 'Produces decisions, action items and consensus from long comment threads'),
]
for i, (num, label, when, desc) in enumerate(items):
    card(s, 0.45, 1.65 + i * 1.02, 12.2, 0.88, fill=CARD_AI, border=BLUE)
    add_rect(s, 0.45, 1.65 + i * 1.02, 0.06, 0.88, BLUE)
    add_text(s, num,   0.65, 1.83 + i * 1.02, 0.65, 0.28, font_size=9,  bold=True, color=BLUE_LIGHT)
    add_text(s, label, 1.45, 1.74 + i * 1.02, 3.1,  0.4,  font_size=13, bold=True)
    add_text(s, when,  1.45, 2.1  + i * 1.02, 3.1,  0.28, font_size=9,  italic=True, color=MUTED)
    divider_v(s, 4.85, 1.74 + i * 1.02, 0.72)
    add_text(s, desc,  5.05, 1.8  + i * 1.02, 7.3,  0.52, font_size=12, color=MUTED)
page_num(s, 11)

# =============================================================================
# SLIDE 12  AI Auto-Fill
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI Feature 1 — Auto-Fill', 'Write a summary. Let AI write the rest.')

card(s, 0.45, 1.18, 5.95, 5.4)
add_text(s, 'How it works', 0.65, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'User types a short issue summary',
    'Clicks Auto-fill with AI',
    'The language model returns:',
], 0.65, 1.82, 5.6, lh=0.52, fs=12)
outputs = ['Full structured description', 'Issue type (Bug, Task, Story ...)', 'Priority level', 'Estimated hours']
for i, o in enumerate(outputs):
    card(s, 0.75, 3.46 + i * 0.72, 5.3, 0.6, fill=NAVY)
    add_text(s, o, 0.95, 3.55 + i * 0.72, 4.9, 0.38, font_size=11)

card(s, 6.75, 1.18, 5.95, 5.4)
add_text(s, 'Key behaviour', 6.95, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'AI never overwrites a field the user has already typed manually',
    'User input always takes priority over AI output',
    'Loading and error states shown during the AI call',
    'Works regardless of issue type or project context',
], 6.95, 1.82, 5.6, lh=0.68, fs=12)
page_num(s, 12)

# =============================================================================
# SLIDE 13  Duplicate Detection
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI Feature 2 — Duplicate Detection', 'Stop creating the same issue twice.')

card(s, 0.45, 1.18, 12.2, 1.45)
add_text(s, 'As the user types a summary, SyncMind runs a real-time similarity check against all existing issues using vector embeddings — a mathematical representation of meaning.',
         0.65, 1.3, 11.8, 0.52, font_size=12)
add_text(s, 'This is not keyword matching. It catches duplicates even when the wording is completely different.',
         0.65, 1.82, 11.8, 0.38, font_size=11, italic=True, color=MUTED)

add_text(s, 'Example', 0.45, 2.85, 2.0, 0.32, font_size=9, bold=True, color=MUTED)
card(s, 0.45, 3.2, 5.6, 0.95)
add_text(s, 'User types:', 0.65, 3.28, 5.2, 0.28, font_size=9, color=MUTED)
add_text(s, '"Login button broken on iPhone"', 0.65, 3.54, 5.2, 0.4, font_size=12, bold=True)

add_text(s, 'Surfaces:', 6.3, 3.28, 1.2, 0.28, font_size=9, color=MUTED)
card(s, 6.6, 3.2, 6.1, 0.95)
add_text(s, '"Authentication fails on mobile Safari"', 6.8, 3.54, 5.7, 0.4, font_size=12, bold=True, color=BLUE_LIGHT)

bullet_list(s, [
    'Debounced  -  runs automatically as the user types, no button needed',
    'Results shown as a warning card below the summary field',
    'Each similar issue is linked so the user can review it before proceeding',
    'Powered by pgvector cosine similarity on stored issue embeddings',
], 0.45, 4.38, 12.2, lh=0.58, fs=12)
page_num(s, 13)

# =============================================================================
# SLIDE 14  Assignee Recommendations
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI Feature 3 — Assignee Recommendations', 'Know who should own the task — and why.')

add_text(s, 'After auto-fill, SyncMind presents a ranked list of team members best suited for the issue. Each suggestion includes a plain-language reason and a one-click Assign action.',
         0.45, 1.15, 12.2, 0.55, font_size=12)

suggestions = [
    ('01', 'Alice Chen',   'Led the authentication refactor last sprint — high familiarity with session handling code', True),
    ('02', 'Bob Martinez', 'Primary mobile engineer, has resolved similar Safari-specific issues previously', False),
    ('03', 'Dana Kim',     'Available capacity this sprint with broad experience across frontend auth flows', False),
]
for i, (rank, name, reason, top) in enumerate(suggestions):
    fill  = CARD_AI if top else CARD
    bord  = BLUE    if top else BORDER
    card(s, 0.45, 1.95 + i * 1.52, 12.2, 1.32, fill=fill, border=bord)
    if top:
        add_rect(s, 0.45, 1.95 + i * 1.52, 0.06, 1.32, BLUE)
    add_text(s, f'Rank {rank}', 0.68, 2.12 + i * 1.52, 0.9, 0.28, font_size=8, bold=True, color=BLUE_LIGHT)
    add_text(s, name,   1.55, 2.06 + i * 1.52, 3.5, 0.44, font_size=14, bold=True)
    add_text(s, reason, 1.55, 2.52 + i * 1.52, 9.4, 0.42, font_size=11, color=MUTED)
    card(s, 11.2, 2.14 + i * 1.52, 1.2, 0.4, fill=BLUE_MID, border=BLUE)
    add_text(s, 'Assign', 11.2, 2.14 + i * 1.52, 1.2, 0.4,
             font_size=10, bold=True, align=PP_ALIGN.CENTER)
page_num(s, 14)

# =============================================================================
# SLIDE 15  Semantic Search
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI Feature 4 — Semantic Search', 'Search by meaning, not just words.')

card(s, 0.45, 1.18, 5.95, 2.35)
add_text(s, 'Keyword Search', 0.65, 1.3, 5.5, 0.4, font_size=13, bold=True, color=MUTED)
add_text(s, 'Standard text matching against issue keys and summaries. Finds issues that contain the exact words typed.',
         0.65, 1.78, 5.5, 1.32, font_size=12, color=MUTED)

card(s, 6.75, 1.18, 5.95, 2.35, fill=CARD_AI, border=BLUE)
top_bar(s, 6.75, 1.18, 5.95)
add_text(s, 'AI Semantic Search', 6.95, 1.3, 5.5, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
add_text(s, 'Query is converted to a vector embedding and compared against every issue using cosine similarity. Finds relevant issues even when no words match.',
         6.95, 1.78, 5.5, 1.32, font_size=12)

card(s, 0.45, 3.72, 12.2, 1.05)
add_text(s, 'EXAMPLE', 0.65, 3.8, 1.5, 0.28, font_size=8, bold=True, color=MUTED)
add_text(s, 'Search:  "user cannot log in"', 0.65, 4.08, 3.8, 0.38, font_size=11, color=MUTED)
divider_v(s, 4.55, 3.9, 0.62)
add_text(s, 'Returns:  "Session expires unexpectedly"  |  "Auth error on login page"  |  "OAuth token refresh fails"',
         4.75, 4.08, 7.7, 0.38, font_size=11, bold=True, color=BLUE_LIGHT)

bullet_list(s, [
    'Toggle between keyword and AI mode from the issue list toolbar',
    'Works together with existing status and priority filters',
    'Powered by pgvector stored on every issue at creation time',
], 0.45, 4.98, 12.2, lh=0.55, fs=12)
page_num(s, 15)

# =============================================================================
# SLIDE 16  Thread Summarisation
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'AI Feature 5 — Thread Summarisation', 'Catch up on any discussion in seconds.')

add_text(s, 'Long comment threads are hard to follow, especially for members joining late. The Summarise Thread button reads all comments and returns a structured breakdown.',
         0.45, 1.15, 12.2, 0.52, font_size=12)

sections = [
    ('Overview',     'What the discussion was about'),
    ('Decisions',    'What was agreed on during the thread'),
    ('Action Items', 'What needs to happen next'),
    ('Consensus',    'The overall conclusion reached'),
]
for i, (label, desc) in enumerate(sections):
    card(s, 0.45 + i * 3.12, 1.88, 2.9, 2.22, fill=CARD_AI, border=BLUE)
    top_bar(s, 0.45 + i * 3.12, 1.88, 2.9)
    add_text(s, label, 0.55 + i * 3.12, 2.08, 2.7, 0.42,
             font_size=13, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, desc,  0.55 + i * 3.12, 2.62, 2.7, 0.8,
             font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

bullet_list(s, [
    'Triggered from the issue slide-over panel — no separate page needed',
    'Summary can be regenerated at any time as the thread grows',
    'A caution note always reminds users that AI output should be verified',
], 0.45, 4.35, 12.2, lh=0.58, fs=12)
page_num(s, 16)

# =============================================================================
# SLIDE 17  How the AI Works
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'How the AI Works', 'Under the hood')

steps = [
    ('01', 'Issue Created',     'Summary and description are sent to the embedding model and stored as a vector in PostgreSQL via pgvector'),
    ('02', 'Similarity Query',  'Semantic search and duplicate detection compare query vectors against stored embeddings using cosine similarity'),
    ('03', 'LLM Call',          'Auto-fill, assignee ranking, and thread summarisation send structured prompts to the language model via OpenRouter'),
    ('04', 'Response Handling', 'Results are parsed, validated, and returned to the frontend — user-touched fields are never overwritten'),
    ('05', 'Single AI Client',  'All AI services resolve through one bound singleton client — making it easy to swap models or providers'),
]
for i, (num, label, desc) in enumerate(steps):
    card(s, 0.45, 1.22 + i * 1.08, 12.2, 0.92)
    add_rect(s, 0.45, 1.22 + i * 1.08, 0.06, 0.92, BLUE)
    add_text(s, num,   0.65, 1.42 + i * 1.08, 0.65, 0.28, font_size=9, bold=True, color=BLUE_LIGHT)
    add_text(s, label, 1.45, 1.3  + i * 1.08, 3.1,  0.42, font_size=12, bold=True)
    divider_v(s, 4.65, 1.35 + i * 1.08, 0.65)
    add_text(s, desc,  4.85, 1.3  + i * 1.08, 7.55, 0.72, font_size=11, color=MUTED)
page_num(s, 17)

# =============================================================================
# SLIDE 18  Future Plans
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Future Plans')

card(s, 0.45, 1.18, 5.95, 5.72)
add_text(s, 'Deeper AI', 0.65, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'AI-generated milestone plans from a goal and deadline',
    'Predictive overdue warnings based on team velocity',
    'Natural language issue creation',
    'AI sprint planning for next milestone backlog',
    'Smart status transitions from comment activity',
], 0.65, 1.82, 5.6, lh=0.6, fs=12)

card(s, 6.75, 1.18, 5.95, 5.72)
add_text(s, 'More Platform Coverage', 6.95, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'Real-time notification centre',
    'Functional global search across all projects',
    'Reporting: burn-down charts, velocity, cycle time',
    'Richer issue metadata: categories, versions, custom fields',
    'User profile data for better AI assignee reasoning',
], 6.95, 1.82, 5.6, lh=0.6, fs=12)
page_num(s, 18)

# =============================================================================
# SLIDE 19  Mobile App
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
slide_title(s, 'Mobile App — Android (Kotlin)', 'The backend is API-first. No server changes needed.')

card(s, 0.45, 1.18, 5.95, 2.52)
add_text(s, 'Approach', 0.65, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'Native Android app built in Kotlin',
    'Calls the same Laravel REST API as the web app',
    'Token-based auth via Laravel Sanctum',
    'Sanctum token support already exists in the backend',
], 0.65, 1.82, 5.6, lh=0.56, fs=12)

card(s, 6.75, 1.18, 5.95, 2.52)
add_text(s, 'Core Mobile Features', 6.95, 1.3, 5.6, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
bullet_list(s, [
    'View and update issues on the go',
    'Post comments and check activity history',
    'Quick milestone progress check',
    'Push notifications for assignments and mentions',
], 6.95, 1.82, 5.6, lh=0.56, fs=12)

card(s, 0.45, 3.92, 12.2, 2.1, fill=CARD_AI, border=BLUE)
add_rect(s, 0.45, 3.92, 0.06, 2.1, BLUE)
add_text(s, 'AI on Mobile', 0.65, 4.02, 11.8, 0.4, font_size=13, bold=True, color=BLUE_LIGHT)
add_text(s, 'All five AI features work through the same API endpoints — auto-fill, duplicate detection, assignee recommendations, semantic search, and thread summarisation are fully available on Android without any additional backend work.',
         0.65, 4.5, 11.8, 0.95, font_size=12)
add_text(s, 'Shared API  |  No duplicate backend effort  |  Consistent AI experience across web and mobile',
         0.45, 6.1, 12.2, 0.32, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)
page_num(s, 19)

# =============================================================================
# SLIDE 20  Closing
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 0.08,  7.5, BLUE)
add_text(s, 'SyncMind', 0.5, 0.85, 12.3, 1.1, font_size=52, bold=True)
add_text(s, 'Built for the way teams actually work.', 0.5, 1.9, 12.3, 0.5,
         font_size=18, italic=True, color=BLUE_LIGHT)
add_rect(s, 0.5, 2.55, 5.5, 0.05, BLUE)

points = [
    'Issue and milestone management that scales with the team',
    'AI that saves real time at every step of the workflow',
    'Clean, fast glassmorphism interface with dark mode',
    'Localized in five languages — English, Myanmar, Japanese, Vietnamese, Khmer',
    'API-first architecture — web today, Android next',
]
for i, p in enumerate(points):
    add_text(s, '  -  ' + p, 0.7, 2.88 + i * 0.63, 12.0, 0.52, font_size=13)

add_rect(s, 0.5, 6.08, 12.3, 0.04, BORDER)
add_text(s, 'Thank you   |   Open for questions', 0.5, 6.28, 12.3, 0.52,
         font_size=18, bold=True, color=BLUE_LIGHT)
page_num(s, 20)

# ── Save ──────────────────────────────────────────────────────────────────────
out = '/tmp/SyncMind_Presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
